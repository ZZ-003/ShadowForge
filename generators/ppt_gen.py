from pptx import Presentation
from pptx.util import Inches, Pt

class PPTGenerator:
    def __init__(self, font_name="Arial", font_size=24):
        self.font_name = font_name
        self.font_size = font_size

    def generate_ppt(self, content_json, output_path, secret=None, secret_placeholder="SECRET_HERE"):
        """
        Generates a PPTX from content.
        content_json: {'slides': [{'title': '...', 'content': ['...']}]}
        """
        
        try:
            prs = Presentation()
            
            slides_data = content_json.get('slides', [])
            
            for slide_data in slides_data:
                title_text = slide_data.get('title', 'Untitled')
                bullets = slide_data.get('content', [])
                
                # Use Title and Content layout (index 1)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                title = slide.shapes.title
                title.text = title_text
                
                # Content
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                
                # First bullet
                if bullets:
                    first_bullet = bullets[0]
                    if secret and secret_placeholder in first_bullet:
                        first_bullet = first_bullet.replace(secret_placeholder, secret)
                    tf.text = first_bullet
                else:
                    tf.text = ""
                
                # Add remaining bullets
                for bullet in bullets[1:]:
                    if secret and secret_placeholder in bullet:
                        bullet = bullet.replace(secret_placeholder, secret)
                    
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            print(f"Error generating PPT: {e}")
            return None
